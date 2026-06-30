"""
Document generation utilities using ReportLab (PDF) and python-pptx (slide decks).
Used for: feasibility reports, proposals, executive summaries (PDF), and pitch decks (PPTX).
"""
import io
import logging

logger = logging.getLogger(__name__)


def generate_feasibility_report_pdf(startup_name: str, content: dict) -> bytes:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=0.75 * inch, bottomMargin=0.75 * inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("TitleCustom", parent=styles["Title"], fontSize=22, spaceAfter=20)
    heading_style = ParagraphStyle("HeadingCustom", parent=styles["Heading2"], spaceBefore=16, spaceAfter=8)
    body_style = styles["BodyText"]

    elements = [
        Paragraph(f"Feasibility Report — {startup_name}", title_style),
        Spacer(1, 12),
    ]

    sections = [
        ("Market Opportunity", content.get("market_opportunity", "Not available.")),
        ("Financial Viability", content.get("financial_viability", "Not available.")),
        ("Operational Readiness", content.get("operational_readiness", "Not available.")),
        ("Key Risks", content.get("key_risks", "Not available.")),
        ("Overall Assessment", content.get("overall_assessment", "Not available.")),
    ]
    for heading, body in sections:
        elements.append(Paragraph(heading, heading_style))
        elements.append(Paragraph(str(body), body_style))

    doc.build(elements)
    return buffer.getvalue()


def generate_proposal_pdf(startup_name: str, content: dict) -> bytes:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=0.75 * inch, bottomMargin=0.75 * inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("TitleCustom", parent=styles["Title"], fontSize=22, spaceAfter=20)
    heading_style = ParagraphStyle("HeadingCustom", parent=styles["Heading2"], spaceBefore=16, spaceAfter=8)
    body_style = styles["BodyText"]

    elements = [
        Paragraph(f"Investment Proposal — {startup_name}", title_style),
        Spacer(1, 12),
    ]

    sections = [
        ("Problem", content.get("problem", "")),
        ("Solution", content.get("solution", "")),
        ("Market", content.get("market", "")),
        ("Team", content.get("team", "")),
        ("Funding Ask & Use of Funds", content.get("funding_ask", "")),
    ]
    for heading, body in sections:
        elements.append(Paragraph(heading, heading_style))
        elements.append(Paragraph(str(body) or "Not provided.", body_style))

    doc.build(elements)
    return buffer.getvalue()


def generate_executive_summary_pdf(startup_name: str, content: dict) -> bytes:
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=1 * inch, bottomMargin=1 * inch)
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle("TitleCustom", parent=styles["Title"], fontSize=20, spaceAfter=16)
    body_style = styles["BodyText"]

    elements = [
        Paragraph(f"Executive Summary — {startup_name}", title_style),
        Spacer(1, 10),
        Paragraph(content.get("summary", "No summary available."), body_style),
    ]
    doc.build(elements)
    return buffer.getvalue()


def generate_pitch_deck_pptx(startup_name: str, content: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    # Title slide
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = startup_name
    slide.placeholders[1].text = content.get("tagline", "Investment Pitch Deck")

    slide_data = [
        ("Problem", content.get("problem", "")),
        ("Solution", content.get("solution", "")),
        ("Market Opportunity", content.get("market", "")),
        ("Team", content.get("team", "")),
        ("The Ask", content.get("funding_ask", "")),
    ]

    for heading, body in slide_data:
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = heading
        body_placeholder = slide.placeholders[1]
        body_placeholder.text = str(body) or "Details to be added."

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def build_document_content(document_type: str, pitch_context: dict) -> dict:
    """
    Maps raw pitch_context (from the conversational flow) into the
    structured content dict expected by each document generator.
    The LLM-driven elaboration (market_opportunity, financial_viability, etc.)
    happens in the agent layer before this — this is the structural assembly step.
    """
    base = {
        "problem": pitch_context.get("problem", ""),
        "solution": pitch_context.get("solution", ""),
        "market": pitch_context.get("market", ""),
        "team": pitch_context.get("team", ""),
        "funding_ask": pitch_context.get("funding_ask", ""),
    }
    if document_type == "executive_summary":
        return {"summary": pitch_context.get("executive_summary_text", "")}
    if document_type == "feasibility_report":
        return {
            "market_opportunity": pitch_context.get("market_opportunity_analysis", base["market"]),
            "financial_viability": pitch_context.get("financial_viability_analysis", ""),
            "operational_readiness": pitch_context.get("operational_readiness_analysis", base["team"]),
            "key_risks": pitch_context.get("key_risks_analysis", ""),
            "overall_assessment": pitch_context.get("overall_assessment", ""),
        }
    return base